from pathlib import Path
from typing import List, Dict, Any
from docx import Document
from pptx import Presentation
import openpyxl
from loguru import logger

from .base import BaseContentProcessor, ProcessedContent, ContentMetadata


class DocumentProcessor(BaseContentProcessor):
    def __init__(self):
        super().__init__()
        self.supported_formats = ['.docx', '.pptx', '.xlsx', '.txt', '.md']
    
    def validate_file(self, file_path: Path) -> bool:
        return file_path.exists() and file_path.suffix.lower() in self.supported_formats
    
    async def process(self, file_path: Path) -> ProcessedContent:
        if not self.validate_file(file_path):
            raise ValueError(f"Invalid document file: {file_path}")
        
        logger.info(f"Processing document: {file_path}")
        
        suffix = file_path.suffix.lower()
        
        if suffix == '.docx':
            return await self._process_docx(file_path)
        elif suffix == '.pptx':
            return await self._process_pptx(file_path)
        elif suffix == '.xlsx':
            return await self._process_xlsx(file_path)
        elif suffix in ['.txt', '.md']:
            return await self._process_text(file_path)
        else:
            raise ValueError(f"Unsupported file format: {suffix}")
    
    async def _process_docx(self, file_path: Path) -> ProcessedContent:
        try:
            doc = Document(file_path)
            
            paragraphs = []
            tables = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append({
                        'text': para.text,
                        'style': para.style.name
                    })
            
            for i, table in enumerate(doc.tables):
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                tables.append({
                    'table_number': i + 1,
                    'data': table_data
                })
            
            full_text = "\n".join([p['text'] for p in paragraphs])
            
            structured_content = {
                'paragraphs': paragraphs,
                'tables': tables,
                'paragraph_count': len(paragraphs),
                'table_count': len(tables)
            }
            
            metadata = self.extract_metadata(file_path)
            metadata.word_count = len(full_text.split())
            metadata.metadata = {
                'paragraph_count': len(paragraphs),
                'table_count': len(tables)
            }
            
            return ProcessedContent(
                raw_text=full_text,
                structured_content=[structured_content],
                metadata=metadata
            )
        except Exception as e:
            logger.error(f"DOCX processing failed: {e}")
            raise
    
    async def _process_pptx(self, file_path: Path) -> ProcessedContent:
        try:
            prs = Presentation(file_path)
            
            slides = []
            full_text = ""
            
            for i, slide in enumerate(prs.slides):
                slide_data = {
                    'slide_number': i + 1,
                    'title': '',
                    'content': [],
                    'notes': ''
                }
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            if shape.is_placeholder and shape.placeholder_format.type == 1:
                                slide_data['title'] = text
                            else:
                                slide_data['content'].append(text)
                            full_text += text + "\n"
                
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text
                    slide_data['notes'] = notes_text
                    full_text += notes_text + "\n"
                
                slides.append(slide_data)
            
            structured_content = {
                'slides': slides,
                'slide_count': len(slides)
            }
            
            metadata = self.extract_metadata(file_path)
            metadata.page_count = len(slides)
            metadata.word_count = len(full_text.split())
            metadata.metadata = {
                'slide_count': len(slides)
            }
            
            return ProcessedContent(
                raw_text=full_text.strip(),
                structured_content=[structured_content],
                metadata=metadata
            )
        except Exception as e:
            logger.error(f"PPTX processing failed: {e}")
            raise
    
    async def _process_xlsx(self, file_path: Path) -> ProcessedContent:
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            
            sheets = []
            full_text = ""
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                
                rows_data = []
                for row in sheet.iter_rows(values_only=True):
                    row_values = [str(cell) if cell is not None else '' for cell in row]
                    if any(row_values):
                        rows_data.append(row_values)
                        full_text += " ".join(row_values) + "\n"
                
                sheets.append({
                    'sheet_name': sheet_name,
                    'data': rows_data,
                    'row_count': len(rows_data),
                    'column_count': sheet.max_column
                })
            
            structured_content = {
                'sheets': sheets,
                'sheet_count': len(sheets)
            }
            
            metadata = self.extract_metadata(file_path)
            metadata.word_count = len(full_text.split())
            metadata.metadata = {
                'sheet_count': len(sheets)
            }
            
            return ProcessedContent(
                raw_text=full_text.strip(),
                structured_content=[structured_content],
                metadata=metadata
            )
        except Exception as e:
            logger.error(f"XLSX processing failed: {e}")
            raise
    
    async def _process_text(self, file_path: Path) -> ProcessedContent:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
            
            lines = text.split('\n')
            structured_content = {
                'lines': [{'line_number': i + 1, 'text': line} for i, line in enumerate(lines)],
                'line_count': len(lines)
            }
            
            metadata = self.extract_metadata(file_path)
            metadata.word_count = len(text.split())
            metadata.metadata = {
                'line_count': len(lines),
                'is_markdown': file_path.suffix.lower() == '.md'
            }
            
            return ProcessedContent(
                raw_text=text,
                structured_content=[structured_content],
                metadata=metadata
            )
        except Exception as e:
            logger.error(f"Text file processing failed: {e}")
            raise
